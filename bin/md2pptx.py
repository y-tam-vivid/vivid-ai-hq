#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Markdown をスライド構成の .pptx へ変換する。

なぜ在るか ── この機には pandoc も LibreOffice も無い（reference_md_to_docx.md）。
企画書・報告書をPowerPointで渡す前提が増えたが、変換の手段が無かった（2026-08-29）。
bin/md2docx.py と同じ思想・同じ検算の作法で作った（新しい流儀を持ち込まない）。

使い方:
    python3 bin/md2pptx.py <入力.md> [出力.pptx]

スライド分割規則（企画書・報告書のMarkdownに出るものだけ。汎用コンバータではない）:
    # 見出し1        → タイトルスライド（表題＋サブタイトルの1枚）
    ## 見出し2       → 新しいスライドの開始（タイトル＝見出し文）
    ### 見出し3      → スライド内の小見出し（太字・やや大きめの本文行）
    | a | b |        → スライド内の表
    - 項目           → 箇条書き（レベル0）
      - 項目         → 箇条書き（インデント2で レベル1）
    1. 項目          → 番号付き（PowerPointには無いので "1. " を残したテキストで表現）
    **強調**         → 太字（行の途中でも効く）
    ---              → 区切り（明示的にスライドを分けたいときの区切り線としては使わない。
                        空段落として本文中に1行差し込む）

★変換した .pptx は必ず PowerPoint（またはKeynote）で開いて目視すること。
  表・箇条書きの階層が意図通りかは、機械では見ていない。
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

JP_FONT = "游ゴシック"

BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


def set_run_font(run, size=None, bold=None):
    run.font.name = JP_FONT
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    # 東アジア文字用フォントも同じ游ゴシックに揃える
    rpr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", JP_FONT)


def add_rich_run(paragraph, text, size=None, base_bold=False):
    """**強調** を太字の run に分けて段落へ流し込む。"""
    pos = 0
    added = False
    for m in BOLD_RE.finditer(text):
        if m.start() > pos:
            r = paragraph.add_run()
            r.text = text[pos:m.start()]
            set_run_font(r, size, base_bold)
            added = True
        r = paragraph.add_run()
        r.text = m.group(1)
        set_run_font(r, size, True)
        added = True
        pos = m.end()
    if pos < len(text):
        r = paragraph.add_run()
        r.text = text[pos:]
        set_run_font(r, size, base_bold)
        added = True
    if not added:
        r = paragraph.add_run()
        r.text = ""
        set_run_font(r, size, base_bold)


def split_row(line):
    return [c.strip() for c in line.strip().strip("|").split("|")]


def is_separator_row(line):
    return bool(re.fullmatch(r"\|[\s:\-|]+\|", line.strip()))


BLANK_LAYOUT_INDEX = 6  # python-pptx既定テンプレの「白紙」レイアウト
TITLE_ONLY_INDEX = 5    # 「タイトルのみ」レイアウト


def new_content_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY_INDEX])
    title_shape = slide.shapes.title
    title_shape.text = title_text
    for p in title_shape.text_frame.paragraphs:
        for r in p.runs:
            set_run_font(r, 28, True)
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(9.0), Inches(5.3))
    tf = body.text_frame
    tf.word_wrap = True
    return slide, tf


def convert(src: Path, dst: Path):
    lines = src.read_text(encoding="utf-8").splitlines()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    cur_tf = None
    cur_first_para_used = False

    def ensure_body():
        nonlocal cur_tf, cur_first_para_used
        if cur_tf is None:
            _, cur_tf = new_content_slide(prs, "")
            cur_first_para_used = False
        return cur_tf

    def next_paragraph(tf):
        nonlocal cur_first_para_used
        if not cur_first_para_used:
            cur_first_para_used = True
            return tf.paragraphs[0]
        return tf.add_paragraph()

    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # 表題スライド（# 見出し1）
        if stripped.startswith("# "):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = stripped[2:]
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    set_run_font(r, 40, True)
            if len(slide.placeholders) > 1:
                sub = slide.placeholders[1]
                sub.text = ""
            cur_tf = None
            cur_first_para_used = False
            i += 1
            continue

        # 新規スライド（## 見出し2）
        if stripped.startswith("## "):
            _, cur_tf = new_content_slide(prs, stripped[3:])
            cur_first_para_used = False
            i += 1
            continue

        # スライド内小見出し（### 見出し3）
        if stripped.startswith("### "):
            tf = ensure_body()
            p = next_paragraph(tf)
            p.space_before = Pt(10)
            add_rich_run(p, stripped[4:], size=18, base_bold=True)
            i += 1
            continue

        # 表
        if stripped.startswith("|"):
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                if not is_separator_row(lines[i]):
                    rows.append(split_row(lines[i]))
                i += 1
            if rows:
                if cur_tf is None:
                    _, cur_tf = new_content_slide(prs, "")
                    cur_first_para_used = False
                slide = prs.slides[-1]
                width = max(len(r) for r in rows)
                rows_n = len(rows)
                left, top = Inches(0.6), Inches(1.6)
                tbl_w, tbl_h = Inches(8.8), Inches(min(0.5 * rows_n, 5.0))
                gtable = slide.shapes.add_table(rows_n, width, left, top, tbl_w, tbl_h)
                table = gtable.table
                for r_idx, row in enumerate(rows):
                    for c_idx in range(width):
                        cell = table.cell(r_idx, c_idx)
                        text = row[c_idx] if c_idx < len(row) else ""
                        cell.text = ""
                        add_rich_run(cell.text_frame.paragraphs[0], text, size=12,
                                      base_bold=(r_idx == 0))
            continue

        # 区切り
        if re.fullmatch(r"-{3,}", stripped):
            tf = ensure_body()
            next_paragraph(tf).text = ""
            i += 1
            continue

        # 箇条書き（インデントで階層を判定。半角2つ or 全角1つを1段とみなす）
        indent_len = len(line) - len(line.lstrip(" "))
        m_ul = re.match(r"^-\s+(.*)$", stripped)
        if m_ul:
            tf = ensure_body()
            p = next_paragraph(tf)
            p.level = min(indent_len // 2, 4)
            add_rich_run(p, m_ul.group(1), size=16)
            i += 1
            continue

        m_ol = re.match(r"^(\d+[.)]\s+.*)$", stripped)
        if m_ol:
            tf = ensure_body()
            p = next_paragraph(tf)
            p.level = min(indent_len // 2, 4)
            add_rich_run(p, m_ol.group(1), size=16)
            i += 1
            continue

        # 引用は本文として扱う
        if stripped.startswith("> "):
            stripped = stripped[2:]

        # 通常の本文
        tf = ensure_body()
        p = next_paragraph(tf)
        add_rich_run(p, stripped, size=16)
        i += 1

    prs.save(str(dst))
    return dst


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    src = Path(sys.argv[1]).expanduser()
    if not src.exists():
        print("入力が見つかりません: %s" % src)
        sys.exit(1)
    dst = Path(sys.argv[2]).expanduser() if len(sys.argv) > 2 else src.with_suffix(".pptx")
    convert(src, dst)
    print("書き出しました: %s (%d バイト)" % (dst, dst.stat().st_size))


if __name__ == "__main__":
    main()
